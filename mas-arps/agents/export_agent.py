import os
import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from state.schema import MASARPSState

def pick_layout(prs: Presentation, slide_data: dict) -> int:
    title = slide_data.get("title", "").lower()
    if slide_data.get("slide_number") == 1:
        return 0   # Title Slide layout
    if "reference" in title:
        return 1   # Title Only layout
    return 1       # Default: Title and Content

def export_node(state: MASARPSState) -> dict:
    slides        = state.get("slides", [])
    template_path = state.get("template_path", "")
    output_path   = state.get("output_path", "output/presentation.pptx")
    now           = datetime.datetime.utcnow().isoformat()

    print(f"[Export] Rendering {len(slides)} slides...")

    # Load template if it exists, otherwise blank presentation
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"[Export] Using template: {template_path}")
    else:
        prs = Presentation()
        print("[Export] No template found — using blank presentation")

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)

    for slide_data in slides:
        layout_idx = pick_layout(prs, slide_data)
        layout     = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
        slide      = prs.slides.add_slide(layout)

        # ── Title ──────────────────────────────────────────
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        # ── Bullets ────────────────────────────────────────
        bullets = slide_data.get("bullets", [])
        body_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_placeholder = ph
                break

        if body_placeholder and bullets:
            tf = body_placeholder.text_frame
            tf.clear()
            for i, bullet_text in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet_text
                p.font.size = Pt(18)
                p.level = 0

            # Add citations below bullets
            citations = slide_data.get("citations", [])
            if citations:
                p = tf.add_paragraph()
                p.text = "  ".join(citations)
                p.font.size = Pt(12)
                p.font.italic = True
                if hasattr(p.font, 'color'):
                    p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        # ── Speaker Notes ──────────────────────────────────
        notes = slide_data.get("speaker_notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    prs.save(output_path)
    print(f"[Export] Saved: {output_path}")

    log_entry = {
        "timestamp": now,
        "node":      "export",
        "action":    "presentation_saved",
        "detail":    output_path
    }

    return {
        "session_log": [*state.get("session_log", []), log_entry]
    }